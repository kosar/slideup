#!/usr/bin/env python3
"""
Enhanced PowerPoint presentation enhancement demo script.

This script demonstrates how to use the PPTXEnhancementCrew to enhance a PowerPoint presentation.
"""

import os
import sys
import argparse
from pathlib import Path

# Add the parent directory to sys.path to enable imports
sys.path.append(str(Path(__file__).resolve().parent.parent.parent))

from agentic.crews.pptx_enhancement_crew import PPTXEnhancementCrew

def main():
    """Run the demo."""
    parser = argparse.ArgumentParser(description='Enhance a PowerPoint presentation')
    parser.add_argument('-i', '--input', type=str, default='tests/demo_files/sample_presentation.pptx',
                        help='Path to the input PowerPoint file')
    parser.add_argument('-o', '--output', type=str, default='tests/output/enhanced_presentation.pptx',
                        help='Path to save the enhanced PowerPoint file')
    parser.add_argument('-v', '--verbose', action='store_true', help='Show detailed output')
    args = parser.parse_args()
    
    # Get API key from environment variable
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        print("Error: OPENAI_API_KEY environment variable not set")
        return 1
    
    # Ensure output directory exists
    output_dir = os.path.dirname(args.output)
    os.makedirs(output_dir, exist_ok=True)
    
    # Initialize the crew with API key
    crew = PPTXEnhancementCrew(api_key=api_key, verbose=args.verbose)
    
    # Enhance the presentation
    try:
        print(f"Enhancing presentation {args.input}...")
        
        # Define enhancement options
        enhancement_options = {
            "improve_design": True,
            "improve_content": True,
            "improve_structure": True,
            "enhancement_level": "moderate",
            "target_audience": "Business professionals",
            "style_guide": {
                "preferred_colors": ["#336699", "#66CCCC", "#FF9966"],
                "preferred_fonts": ["Arial", "Georgia"]
            }
        }
        
        # Use the enhance_presentation method
        result = crew.enhance_presentation(
            input_file=args.input, 
            output_file=args.output,
            enhancement_options=enhancement_options
        )
        
        print(f"Enhancement completed successfully!")
        print(f"Enhanced presentation saved to: {args.output}")
        return 0
        
    except Exception as e:
        print(f"Error: {e}")
        return 1

def create_sample_presentation():
    """Create a sample PowerPoint presentation for testing if needed."""
    from pptx import Presentation
    
    # This is a simplified version that just creates a basic presentation
    prs = Presentation()
    
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sample Presentation"
    subtitle.text = "For Testing Enhancement"
    
    # Add a content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Bullet Points"
    content.text = "• First bullet point\n• Second bullet point\n• Third bullet point"
    
    # Ensure directory exists
    os.makedirs("tests/demo_files", exist_ok=True)
    
    # Save the presentation
    file_path = "tests/demo_files/sample_presentation.pptx"
    prs.save(file_path)
    print(f"Created sample presentation at: {file_path}")
    
    return file_path

if __name__ == "__main__":
    # Check if the sample file exists, create it if not
    sample_file = Path("tests/demo_files/sample_presentation.pptx")
    if not sample_file.exists():
        create_sample_presentation()
        
    exit(main())
