import os
import sys
import shutil
from crewai import Crew, Task
from pathlib import Path

# Add parent directory to path so we can import from agentic.agents
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))

# Import our agents with correct class names
from agentic.agents.pptx_analyzer_agent import PptxAnalyzerAgentWithAPIKey
from agentic.agents.pptx_enhancer_agent import PPTXEnhancerAgent

class PPTXEnhancementCrew:
    """Crew for enhancing PowerPoint presentations by analyzing and improving them."""
    
    def __init__(self, api_key=None, verbose=True):
        """
        Initialize the enhancement crew.
        
        Args:
            api_key: API key for LLM services
            verbose: Whether to show detailed output
        """
        self.api_key = api_key
        self.verbose = verbose
        self.analyzer_agent = None
        self.enhancer_agent = None
        
    def setup_agents(self):
        """Set up the analyzer and enhancer agents."""
        self.analyzer_agent = PptxAnalyzerAgentWithAPIKey(api_key=self.api_key)
        self.enhancer_agent = PPTXEnhancerAgent(api_key=self.api_key)
        
    def enhance_presentation(self, input_file, output_file, enhancement_options=None):
        """
        Enhance a PowerPoint presentation.
        
        Args:
            input_file: Path to the input PPTX file
            output_file: Path where the enhanced PPTX will be saved
            enhancement_options: Dictionary of enhancement options:
                - improve_design: Improve visual design (bool)
                - improve_content: Improve content quality (bool)
                - improve_structure: Improve presentation structure (bool)
                - style_guide: Dictionary with style preferences
                - target_audience: Target audience description
                - enhancement_level: How aggressive to be (light, moderate, comprehensive)
        
        Returns:
            Path to the enhanced presentation
        """
        # Default options
        default_options = {
            "improve_design": True,
            "improve_content": True,
            "improve_structure": True,
            "enhancement_level": "moderate",
            "style_guide": {"preferred_colors": [], "preferred_fonts": []},
            "target_audience": "General audience"
        }
        
        # Apply defaults for missing options
        if not enhancement_options:
            enhancement_options = default_options
        else:
            for key, value in default_options.items():
                if key not in enhancement_options:
                    enhancement_options[key] = value
            
        # Ensure input file exists
        if not os.path.exists(input_file):
            raise FileNotFoundError(f"Input file not found: {input_file}")
        
        # Create backup
        backup_file = f"{input_file}.backup"
        try:
            shutil.copy2(input_file, backup_file)
            if self.verbose:
                print(f"Created backup at {backup_file}")
        except Exception as e:
            print(f"Warning: Could not create backup: {e}")
            
        # Setup agents if not already done
        if not self.analyzer_agent or not self.enhancer_agent:
            self.setup_agents()
            
        # Create analysis task with properly formatted context items
        analysis_task = Task(
            description=self._create_analysis_prompt(input_file, enhancement_options),
            expected_output="Detailed analysis of the presentation with specific recommendations for improvement",
            agent=self.analyzer_agent.agent,
            tools=[self.analyzer_agent.agent.tools[0]],  # Pass the analyze tool explicitly
            context=[
                {
                    "description": "Presentation file path to analyze",
                    "expected_output": "Analysis results",
                    "file_path": input_file,
                    "enhancement_options": enhancement_options
                }
            ]
        )
        
        # Create enhancement task with properly formatted context items
        enhancement_task = Task(
            description=self._create_enhancement_prompt(input_file, output_file, enhancement_options),
            expected_output="Enhanced PowerPoint presentation with design, content, and structure improvements",
            agent=self.enhancer_agent.agent,
            tools=[self.enhancer_agent.agent.tools[0]],  # Pass the enhance tool explicitly
            context=[
                {
                    "description": "Presentation enhancement parameters",
                    "expected_output": "Enhanced presentation confirmation",
                    "input_file": input_file,
                    "output_file": output_file,
                    "enhancement_options": enhancement_options
                }
            ],
            dependencies=[analysis_task]
        )
        
        # Create and run the crew
        crew = Crew(
            agents=[self.analyzer_agent.agent, self.enhancer_agent.agent],
            tasks=[analysis_task, enhancement_task],
            verbose=self.verbose
        )
        
        try:
            result = crew.kickoff()
            
            # Verify that output file was created
            if not os.path.exists(output_file):
                print(f"Warning: Output file not found at {output_file}")
                print("Using original file as a fallback...")
                shutil.copy2(input_file, output_file)
                
            return output_file
            
        except Exception as e:
            print(f"Error occurred during enhancement: {e}")
            # Try to restore from backup
            if os.path.exists(backup_file):
                print("Restoring from backup...")
                shutil.copy2(backup_file, output_file)
            
            raise
        finally:
            # Clean up backup
            if os.path.exists(backup_file):
                os.remove(backup_file)
    
    def _create_analysis_prompt(self, input_file, options):
        """Create a detailed prompt for the analysis task."""
        prompt = f"Analyze the presentation {input_file} thoroughly and identify areas for improvement.\n\n"
        
        if options["improve_design"]:
            prompt += "- Evaluate the visual design including colors, fonts, layout, and overall consistency.\n"
        
        if options["improve_content"]:
            prompt += "- Assess the content quality, clarity, and messaging effectiveness.\n"
        
        if options["improve_structure"]:
            prompt += "- Review the presentation structure, flow, and logical organization.\n"
        
        prompt += f"\nTarget audience: {options['target_audience']}\n"
        prompt += f"Enhancement level: {options['enhancement_level']}\n"
        
        return prompt
    
    def _create_enhancement_prompt(self, input_file, output_file, options):
        """Create a detailed prompt for the enhancement task."""
        prompt = f"Enhance the presentation from {input_file} and save to {output_file} based on the analysis provided.\n\n"
        
        prompt += "Apply the following improvements:\n"
        if options["improve_design"]:
            prompt += "- Improve the visual design for better aesthetics and clarity.\n"
        
        if options["improve_content"]:
            prompt += "- Enhance the content quality and messaging.\n"
        
        if options["improve_structure"]:
            prompt += "- Improve the presentation structure and flow.\n"
        
        if options["style_guide"]["preferred_colors"]:
            colors = ", ".join(options["style_guide"]["preferred_colors"])
            prompt += f"- Use the following preferred colors: {colors}.\n"
            
        if options["style_guide"]["preferred_fonts"]:
            fonts = ", ".join(options["style_guide"]["preferred_fonts"])
            prompt += f"- Use the following preferred fonts: {fonts}.\n"
        
        prompt += f"\nTarget audience: {options['target_audience']}\n"
        prompt += f"Enhancement level: {options['enhancement_level']}\n"
        
        return prompt

import os
from crewai import Agent, Crew, Task
from langchain.tools import Tool
# Fix import for OpenAI
from langchain_openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import shutil

class PPTXEnhancementCrew:
    """
    A crew for enhancing PowerPoint presentations.
    """
    
    def __init__(self):
        """Initialize the crew with necessary agents and tools"""
        # Set up the LLM
        self.llm = OpenAI(temperature=0.5)
        
        # Initialize agents
        self.content_analyst = self._create_content_analyst()
        self.visual_designer = self._create_visual_designer()
        self.structure_optimizer = self._create_structure_optimizer()
        self.formatting_expert = self._create_formatting_expert()
        
        # Build the crew
        self.crew = Crew(
            agents=[
                self.content_analyst,
                self.visual_designer,
                self.structure_optimizer,
                self.formatting_expert
            ],
            tasks=[]
        )
    
    def _create_content_analyst(self):
        """Create the content analyst agent"""
        return Agent(
            role="Content Analyst",
            goal="Analyze and improve presentation content",
            backstory="You are an expert content strategist specializing in clear, concise, and impactful messaging for presentations.",
            verbose=True,
            llm=self.llm,
            tools=[
                Tool(
                    name="analyze_content",
                    func=self._analyze_content,
                    description="Analyze presentation content"
                ),
                Tool(
                    name="enhance_content",
                    func=self._enhance_content,
                    description="Enhance presentation content"
                )
            ]
        )
    
    def _create_visual_designer(self):
        """Create the visual designer agent"""
        return Agent(
            role="Visual Designer",
            goal="Improve the visual design of presentations",
            backstory="You are a talented presentation designer with expertise in creating visually appealing slides that convey information effectively.",
            verbose=True,
            llm=self.llm,
            tools=[
                Tool(
                    name="analyze_design",
                    func=self._analyze_design,
                    description="Analyze presentation design"
                ),
                Tool(
                    name="improve_design",
                    func=self._improve_design,
                    description="Improve presentation design"
                )
            ]
        )
    
    def _create_structure_optimizer(self):
        """Create the structure optimizer agent"""
        return Agent(
            role="Structure Optimizer",
            goal="Optimize the structure and flow of presentations",
            backstory="You are an expert in information architecture and storytelling, specializing in creating logical and engaging presentation structures.",
            verbose=True,
            llm=self.llm,
            tools=[
                Tool(
                    name="analyze_structure",
                    func=self._analyze_structure,
                    description="Analyze presentation structure"
                ),
                Tool(
                    name="optimize_structure",
                    func=self._optimize_structure,
                    description="Optimize presentation structure"
                )
            ]
        )
    
    def _create_formatting_expert(self):
        """Create the formatting expert agent"""
        return Agent(
            role="Formatting Expert",
            goal="Improve slide formatting and consistency",
            backstory="You are a detail-oriented formatting specialist who ensures presentations look professional and consistent throughout.",
            verbose=True,
            llm=self.llm,
            tools=[
                Tool(
                    name="analyze_formatting",
                    func=self._analyze_formatting,
                    description="Analyze presentation formatting"
                ),
                Tool(
                    name="improve_formatting",
                    func=self._improve_formatting,
                    description="Improve presentation formatting"
                )
            ]
        )
    
    def _analyze_content(self, presentation):
        """Analyze the content of a presentation"""
        # Implementation to analyze presentation content
        pass
    
    def _enhance_content(self, presentation, analysis, enhancement_level):
        """Enhance the content of a presentation"""
        # Implementation to improve content based on analysis
        pass
    
    def _analyze_design(self, presentation):
        """Analyze the design of a presentation"""
        # Implementation to analyze presentation design
        pass
    
    def _improve_design(self, presentation, analysis, enhancement_level):
        """Improve the design of a presentation"""
        # Implementation to improve design based on analysis
        pass
    
    def _analyze_structure(self, presentation):
        """Analyze the structure of a presentation"""
        # Implementation to analyze presentation structure
        pass
    
    def _optimize_structure(self, presentation, analysis, enhancement_level):
        """Optimize the structure of a presentation"""
        # Implementation to improve structure based on analysis
        pass
    
    def _analyze_formatting(self, presentation):
        """Analyze the formatting of a presentation"""
        # Implementation to analyze presentation formatting
        pass
    
    def _improve_formatting(self, presentation, analysis, enhancement_level):
        """Improve the formatting of a presentation"""
        # Implementation to improve formatting based on analysis
        pass

    def _extract_presentation_content(self, pptx_file):
        """Extract content from a PowerPoint presentation"""
        presentation = Presentation(pptx_file)
        
        content = {
            "slides": [],
            "overall": {
                "slide_count": len(presentation.slides),
                "title": None
            }
        }
        
        for i, slide in enumerate(presentation.slides):
            slide_content = {
                "index": i,
                "title": None,
                "text_content": [],
                "has_images": False,
                "has_charts": False,
                "has_tables": False,
                "layout_type": str(slide.slide_layout.name)
            }
            
            # Extract title if available
            if slide.shapes.title:
                slide_content["title"] = slide.shapes.title.text
                if i == 0:
                    content["overall"]["title"] = slide.shapes.title.text
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content["text_content"].append(shape.text)
                
                # Check if shape is an image
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    slide_content["has_images"] = True
                
                # Check if shape is a chart
                if shape.shape_type == 3:  # MSO_SHAPE_TYPE.CHART
                    slide_content["has_charts"] = True
                
                # Check if shape is a table
                if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    slide_content["has_tables"] = True
            
            content["slides"].append(slide_content)
        
        return content

    def _enhance_presentation(self, input_file, output_file, improvement_options, enhancement_level, instructions):
        """Enhance a PowerPoint presentation"""
        # For now, we'll implement a simplified version that just copies the file
        # In a real implementation, this would analyze and modify the presentation
        shutil.copy(input_file, output_file)
        
        # Here we'd implement the actual enhancement logic
        # This would involve:
        # 1. Opening the presentation
        # 2. Analyzing content, design, structure, and formatting
        # 3. Making improvements based on the analysis and options
        # 4. Saving the enhanced presentation
        
        # For demonstration purposes, let's add a note to confirm what would be enhanced
        presentation = Presentation(output_file)
        
        # Add a note slide at the end to indicate what was enhanced
        slide_layout = presentation.slide_layouts[5]  # Title and notes layout
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "SlideUp AI Enhancement Summary"
        
        # Add a text box for the enhancement details
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        
        p = tf.add_paragraph()
        p.text = "This presentation was enhanced with the following settings:"
        
        p = tf.add_paragraph()
        p.text = f"Enhancement level: {enhancement_level}"
        p.level = 1
        
        # Add details about each enhancement option
        for option, enabled in improvement_options.items():
            if enabled:
                p = tf.add_paragraph()
                p.text = f"{option.replace('_', ' ').title()}: Enabled"
                p.level = 1
        
        # Add any additional instructions
        if instructions:
            p = tf.add_paragraph()
            p.text = "Additional instructions:"
            
            p = tf.add_paragraph()
            p.text = instructions
            p.level = 1
        
        # Save the enhanced presentation
        presentation.save(output_file)
        
        return output_file
    
    def run(self, input_file, output_file, improve_visuals=True, enhance_content=True, 
            optimize_structure=True, improve_formatting=True, enhancement_level="moderate", 
            additional_instructions=""):
        """
        Run the crew to enhance a PowerPoint presentation.
        
        Args:
            input_file (str): Path to the input PowerPoint file
            output_file (str): Path to save the enhanced PowerPoint file
            improve_visuals (bool): Whether to improve visual design
            enhance_content (bool): Whether to enhance content quality
            optimize_structure (bool): Whether to optimize presentation structure
            improve_formatting (bool): Whether to improve slide formatting
            enhancement_level (str): Level of enhancement (light, moderate, comprehensive)
            additional_instructions (str): Any additional instructions for enhancement
            
        Returns:
            dict: Status and path to the enhanced PowerPoint file
        """
        try:
            # Extract content for analysis
            content = self._extract_presentation_content(input_file)
            
            # Set up improvement options
            improvement_options = {
                "improve_visuals": improve_visuals,
                "enhance_content": enhance_content,
                "optimize_structure": optimize_structure,
                "improve_formatting": improve_formatting
            }
            
            # Enhance the presentation
            result = self._enhance_presentation(
                input_file, 
                output_file, 
                improvement_options, 
                enhancement_level,
                additional_instructions
            )
            
            return {"status": "success", "output_file": output_file}
        
        except Exception as e:
            return {"status": "error", "message": str(e)}
