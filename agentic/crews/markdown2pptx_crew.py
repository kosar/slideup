import os
import sys
from typing import Optional
import json

from crewai import Crew, Task, Agent

# Add the project root to the Python path
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
if project_root not in sys.path:
    sys.path.append(project_root)

# Use relative imports
from ..agents.markdown_analyzer import MarkdownAnalyzerAgent
from ..agents.pptx_generator import PPTXGeneratorAgent
from ..utils.logger import setup_logger

logger = setup_logger("MarkdownToPPTXCrew")

class MarkdownToPPTXCrew:
    def __init__(self, use_direct_calls=False):
        self.use_direct_calls = use_direct_calls

    def run(self, markdown_text, output_path):
        # Placeholder for the actual conversion logic
        with open(output_path, "w") as f:
            f.write("This is a placeholder for the converted PowerPoint file.")
        return output_path

def main():
    """
    A simple demonstration of using the MarkdownToPPTXCrew independently of a web interface.
    """
    # Sample Markdown content
    sample_markdown = """
    # Title Slide
    ## Presenter: John Doe
    ## Date: 2024-01-01

    # First Content Slide
    - Introduction
    - Main point
    - Conclusion

    ```python
    def hello_world():
        print("Hello, world!")
    ```
    """

    # Define the output file path
    output_file_path = "output_presentation.pptx"

    # Create an instance of the MarkdownToPPTXCrew
    markdown_to_pptx = MarkdownToPPTXCrew(use_direct_calls=True)

    # Run the crew to convert the Markdown content to a PowerPoint presentation
    result = markdown_to_pptx.run(sample_markdown, output_file_path)

    if result:
        print(f"PowerPoint presentation created successfully at: {result}")
    else:
        print("Failed to create PowerPoint presentation.")

if __name__ == "__main__":
    main()

import os
from crewai import Agent, Crew, Task
from langchain.tools import Tool
# Fix import for OpenAI
from langchain_openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from markdown import Markdown
from bs4 import BeautifulSoup
import re

class Markdown2PPTXCrew:
    """
    A crew for converting Markdown files to PowerPoint presentations.
    """
    
    def __init__(self):
        """Initialize the crew with necessary agents and tools"""
        # Set up the LLM
        self.llm = OpenAI(temperature=0.7)
        
        # Initialize agents
        self.content_analyzer = self._create_content_analyzer()
        self.slide_designer = self._create_slide_designer()
        self.presentation_builder = self._create_presentation_builder()
        
        # Build the crew
        self.crew = Crew(
            agents=[
                self.content_analyzer,
                self.slide_designer,
                self.presentation_builder
            ],
            tasks=[]
        )
    
    def _create_content_analyzer(self):
        """Create the content analyzer agent"""
        return Agent(
            role="Content Analyzer",
            goal="Analyze markdown content and extract structured information for slides",
            backstory="You are an expert in content analysis and information architecture, specializing in identifying key points and organizing them into presentation structures.",
            verbose=True,
            llm=self.llm,
            tools=[
                Tool(
                    name="parse_markdown",
                    func=self._parse_markdown,
                    description="Parse markdown text into structured content"
                )
            ]
        )
    
    def _create_slide_designer(self):
        """Create the slide designer agent"""
        return Agent(
            role="Slide Designer",
            goal="Design visually appealing slides based on the content structure",
            backstory="You are a professional presentation designer with a keen eye for layout, typography, and visual hierarchy.",
            verbose=True,
            llm=self.llm,
            tools=[
                Tool(
                    name="generate_slide_layouts",
                    func=self._generate_slide_layouts,
                    description="Generate slide layouts based on content type and theme"
                )
            ]
        )
    
    def _create_presentation_builder(self):
        """Create the presentation builder agent"""
        return Agent(
            role="Presentation Builder",
            goal="Build a PowerPoint presentation by combining content and design",
            backstory="You are an expert in creating PowerPoint presentations that effectively communicate key messages with optimal visual design.",
            verbose=True,
            llm=self.llm,
            tools=[
                Tool(
                    name="build_presentation",
                    func=self._build_presentation,
                    description="Build a PowerPoint presentation"
                )
            ]
        )
    
    def _parse_markdown(self, markdown_text):
        """Parse markdown text into structured content"""
        # Implementation to parse markdown and identify sections, lists, etc.
        # This is called via the content_analyzer agent
        pass
    
    def _generate_slide_layouts(self, content_structure, theme, color_scheme):
        """Generate slide layouts based on content type and theme"""
        # Implementation to create slide templates and styles
        # This is called via the slide_designer agent
        pass
    
    def _build_presentation(self, structured_content, slide_layouts, output_file):
        """Build a PowerPoint presentation based on structured content and layouts"""
        # Implementation to create the actual PowerPoint presentation
        # This is called via the presentation_builder agent
        pass
        
    def _extract_sections_from_markdown(self, markdown_text):
        """Extract sections, subsections and content from markdown text"""
        # Convert markdown to HTML
        html = Markdown().convert(markdown_text)
        soup = BeautifulSoup(html, 'html.parser')
        
        # Extract sections by headers
        sections = []
        current_section = {"title": None, "content": [], "level": 0}
        
        for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol']):
            if element.name.startswith('h'):
                level = int(element.name[1])
                
                if current_section["title"]:
                    sections.append(current_section)
                
                current_section = {
                    "title": element.text.strip(),
                    "content": [],
                    "level": level
                }
            else:
                current_section["content"].append({
                    "type": element.name,
                    "content": str(element)
                })
        
        if current_section["title"]:
            sections.append(current_section)
            
        return sections
    
    def _create_pptx(self, sections, output_file, theme, color_scheme, include_cover, include_agenda):
        """Create a PowerPoint presentation from the parsed sections"""
        # Create presentation
        prs = Presentation()
        
        # Apply theme and color scheme
        theme_settings = self._get_theme_settings(theme, color_scheme)
        
        # Create title slide if requested
        if include_cover:
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            if sections and sections[0]["level"] == 1:
                title.text = sections[0]["title"]
            else:
                title.text = "Presentation"
                
            subtitle.text = "Generated with SlideUp AI"
            
        # Create agenda slide if requested
        if include_agenda:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Agenda"
            
            content = slide.placeholders[1]
            tf = content.text_frame
            
            for section in sections:
                if section["level"] <= 2:  # Only include top-level sections in agenda
                    p = tf.add_paragraph()
                    p.text = section["title"]
                    p.level = section["level"]
        
        # Create content slides
        for section in sections:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = section["title"]
            
            # Add content
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            
            for item in section["content"]:
                p = tf.add_paragraph()
                
                if item["type"] == "p":
                    # Regular paragraph
                    soup = BeautifulSoup(item["content"], 'html.parser')
                    p.text = soup.get_text()
                elif item["type"] in ["ul", "ol"]:
                    # List items
                    soup = BeautifulSoup(item["content"], 'html.parser')
                    for i, li in enumerate(soup.find_all('li')):
                        if i > 0:
                            p = tf.add_paragraph()
                        p.text = li.get_text()
                        p.level = 1  # Indent list items
        
        # Save the presentation
        prs.save(output_file)
        
        return output_file
    
    def _get_theme_settings(self, theme, color_scheme):
        """Get theme and color scheme settings"""
        themes = {
            "Professional": {
                "Blue": {"primary": "0x1F497D", "secondary": "0xCCDFFF"},
                "Green": {"primary": "0x385E0F", "secondary": "0xE3F0D7"},
                "Red": {"primary": "0x982B2B", "secondary": "0xF5D7D7"},
                "Purple": {"primary": "0x5F497A", "secondary": "0xE4DFEC"},
                "Orange": {"primary": "0xE46C0A", "secondary": "0xFDE9D9"},
                "Grayscale": {"primary": "0x404040", "secondary": "0xF2F2F2"},
            },
            "Creative": {
                # Add color schemes for Creative theme
            },
            "Minimalist": {
                # Add color schemes for Minimalist theme
            },
            "Academic": {
                # Add color schemes for Academic theme
            },
            "Corporate": {
                # Add color schemes for Corporate theme
            }
        }
        
        return themes.get(theme, {}).get(color_scheme, {"primary": "0x1F497D", "secondary": "0xCCDFFF"})
    
    def run(self, markdown_file, output_file, theme="Professional", color_scheme="Blue", include_cover=True, include_agenda=True):
        """
        Run the crew to convert a markdown file to a PowerPoint presentation.
        
        Args:
            markdown_file (str): Path to the markdown file
            output_file (str): Path to save the PowerPoint file
            theme (str): Presentation theme
            color_scheme (str): Color scheme
            include_cover (bool): Whether to include a cover slide
            include_agenda (bool): Whether to include an agenda slide
            
        Returns:
            str: Path to the generated PowerPoint file
        """
        try:
            # Read markdown file
            with open(markdown_file, 'r', encoding='utf-8') as f:
                markdown_text = f.read()
            
            # Parse markdown into sections
            sections = self._extract_sections_from_markdown(markdown_text)
            
            # Create PowerPoint presentation
            result = self._create_pptx(sections, output_file, theme, color_scheme, include_cover, include_agenda)
            
            return {"status": "success", "output_file": output_file}
        
        except Exception as e:
            return {"status": "error", "message": str(e)}
