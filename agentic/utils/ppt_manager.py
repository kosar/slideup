"""
PowerPoint manipulation base class using python-pptx.
"""

from pathlib import Path
from typing import Union, List, Tuple, Optional, Dict, Any
import os

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.autoshape import Shape
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PPTManager:
    """Base class for PowerPoint presentation manipulation."""
    
    def __init__(self, template_path: Optional[Union[str, Path]] = None):
        """
        Initialize the PowerPoint manager.
        
        Args:
            template_path: Path to a template PPTX file. If None, creates a blank presentation.
        """
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
    def save(self, file_path: Union[str, Path]) -> None:
        """
        Save the presentation to a file.
        
        Args:
            file_path: Path where the presentation will be saved
        """
        Path(file_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(file_path)
    
    def add_slide(self, layout_index: int = 1) -> Slide:
        """
        Add a new slide to the presentation.
        
        Args:
            layout_index: Index of the slide layout to use
            
        Returns:
            The newly created slide
        """
        layout = self.prs.slide_layouts[layout_index]
        return self.prs.slides.add_slide(layout)
    
    def add_title_slide(self, title: str, subtitle: str = "") -> Slide:
        """
        Add a title slide to the presentation.
        
        Args:
            title: Title text
            subtitle: Subtitle text
            
        Returns:
            The newly created title slide
        """
        slide = self.add_slide(0)  # 0 is usually the title slide layout
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        if title_shape:
            title_shape.text = title
        if subtitle_shape:
            subtitle_shape.text = subtitle
            
        return slide
    
    def add_content_slide(self, title: str, content_items: List[str]) -> Slide:
        """
        Add a content slide with bullet points.
        
        Args:
            title: Slide title
            content_items: List of bullet points
            
        Returns:
            The newly created content slide
        """
        slide = self.add_slide(1)  # 1 is usually a title and content layout
        
        # Set the title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Add bullet points to the content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # 2 is for body/content
                content_placeholder = shape
                break
        
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for i, item in enumerate(content_items):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = item
                p.level = 0  # Top-level bullet
        
        return slide
    
    def add_section(self, title: str) -> None:
        """Add a new section to the presentation."""
        # Get current index of slides
        idx = len(self.prs.slides) - 1 if self.prs.slides else -1
        self.prs.sections.add_section(title, idx + 1)
    
    def add_picture(self, slide: Slide, image_path: Union[str, Path], 
                   left: float = 1.0, top: float = 2.0, 
                   width: Optional[float] = None, height: Optional[float] = None) -> Shape:
        """
        Add a picture to a slide.
        
        Args:
            slide: Slide to add the picture to
            image_path: Path to the image file
            left: Left position in inches
            top: Top position in inches
            width: Width in inches (aspect ratio maintained if only one dimension provided)
            height: Height in inches
            
        Returns:
            The picture shape
        """
        return slide.shapes.add_picture(
            str(image_path), 
            Inches(left), 
            Inches(top),
            Inches(width) if width else None,
            Inches(height) if height else None
        )
    
    def add_textbox(self, slide: Slide, text: str, 
                   left: float = 1.0, top: float = 2.0,
                   width: float = 4.0, height: float = 1.0,
                   font_size: int = 18,
                   alignment: PP_ALIGN = PP_ALIGN.LEFT) -> Shape:
        """
        Add a textbox to a slide.
        
        Args:
            slide: Slide to add the textbox to
            text: Text content
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            font_size: Font size in points
            alignment: Text alignment
            
        Returns:
            The textbox shape
        """
        textbox = slide.shapes.add_textbox(
            Inches(left), 
            Inches(top),
            Inches(width),
            Inches(height)
        )
        
        tf = textbox.text_frame
        tf.text = text
        
        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(font_size)
        paragraph.alignment = alignment
        
        return textbox
    
    def apply_template(self, template_path: Union[str, Path]) -> None:
        """
        Apply a template to the current presentation.
        
        Args:
            template_path: Path to the template PPTX file
        """
        # This is a simplified implementation - a full implementation would
        # copy slide masters and layouts from the template
        template = Presentation(template_path)
        
        # Copy slide masters and layouts
        for master in template.slide_masters:
            for layout in master.slide_layouts:
                # This is a placeholder - in reality, copying layouts requires
                # more complex XML manipulation
                pass
    
    def add_chart(self, slide: Slide, chart_type: str, data: Dict[str, List],
                 left: float = 1.0, top: float = 2.0,
                 width: float = 6.0, height: float = 4.0) -> None:
        """
        Add a chart to a slide (placeholder implementation).
        
        Args:
            slide: Slide to add the chart to
            chart_type: Type of chart to add
            data: Chart data
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
        """
        # This is a placeholder - implementing full chart functionality
        # requires more detailed code with python-pptx
        chart_placeholder = slide.shapes.add_chart(
            chart_type, 
            Inches(left), 
            Inches(top),
            Inches(width),
            Inches(height)
        )
        # In a real implementation, you would populate the chart data here
