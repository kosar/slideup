from pptx import Presentation
import os

# Import existing tool if available; otherwise, re-implement minimal version
try:
    from pptx2enhancedpptx import extract_slide_content
except ImportError:
    # Minimal version of extract_slide_content if import fails
    def extract_slide_content(prs):
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    slide_text.append(shape.text.strip())
            slides_text.append({'index': i, 'text': " ".join(slide_text).strip()})
        return slides_text

def analyze_presentation(pptx_path):
    # Load presentation and extract slide contents
    if not os.path.exists(pptx_path):
        return f"Error: File '{pptx_path}' does not exist."
    
    prs = Presentation(pptx_path)
    slides_text = extract_slide_content(prs)
    
    report_lines = []
    total_slides = len(prs.slides)
    report_lines.append("Presentation Analysis Report")
    report_lines.append(f"Total slides: {total_slides}\n")
    
    # Analyze each slide
    for slide in prs.slides:
        index = prs.slides.index(slide)
        slide_info = next((item for item in slides_text if item['index'] == index), {})
        text_content = slide_info.get('text', "")
        issues = []
        
        # Check for empty or low content
        if not text_content:
            issues.append("No text content found.")
        elif len(text_content) < 20:
            issues.append("Possibly sparse text; consider adding more details.")
        elif len(text_content) > 500:
            issues.append("Content appears crowded; consider splitting or summarizing.")
        
        # Check if slide has any pictures
        picture_found = any(hasattr(shape, "shape_type") and shape.shape_type == 13 for shape in slide.shapes)
        if not picture_found:
            issues.append("No images detected; adding relevant visuals may enhance the design.")
        
        # Check for inconsistent font sizes in text shapes
        font_sizes = set()
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font.size:
                        font_sizes.add(para.font.size.pt)
        if len(font_sizes) > 1:
            issues.append("Inconsistent font sizes detected; standardize for visual clarity.")
        
        # -- New Check: Inconsistent font families --
        font_names = set()
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            font_names.add(run.font.name)
        if len(font_names) > 1:
            issues.append("Multiple font families detected; standardize for cohesion.")
        
        # -- New Check: Inconsistent text alignments --
        alignments = set()
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.alignment is not None:
                        alignments.add(para.alignment)
        if len(alignments) > 1:
            issues.append("Inconsistent text alignments detected; consider standardizing alignment.")
        
        # Compile slide analysis
        report_lines.append(f"Slide {index + 1}:")
        if issues:
            for issue in issues:
                report_lines.append(f"  - {issue}")
        else:
            report_lines.append("  - No major issues detected.")
        report_lines.append("")  # Empty line between slides
        
    # Final recommendations
    report_lines.append("General Recommendations:")
    report_lines.append("  - Ensure slide layouts are consistent across the presentation.")
    report_lines.append("  - Use visuals to break up text-heavy slides.")
    report_lines.append("  - Review typography and spacing for enhanced readability.")
    report_lines.append("  - Standardize font families and text alignments for a cohesive design.")
    
    return "\n".join(report_lines)

# CrewAI agent class for analyzing PPTX files
class PptxAnalyzerAgent:
    def __init__(self):
        # ...any initialization if needed...
        pass

    def run(self, pptx_path):
        """
        Accepts PPTX file path and returns an analysis report.
        """
        return analyze_presentation(pptx_path)

    def process(self, request):
        """
        Process a CrewAI request. Expects 'pptx_path' in the request dict.
        """
        pptx_path = request.get('pptx_path')
        if not pptx_path:
            return {"error": "No PPTX file path provided in request."}
        report = self.run(pptx_path)
        return {"analysis_report": report}

if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print("Usage: python pptx_analyzer.py <presentation_file.pptx>")
    else:
        pptx_path = sys.argv[1]
        agent = PptxAnalyzerAgent()
        result = agent.process({'pptx_path': pptx_path})
        if 'error' in result:
            print(result['error'])
        else:
            print(result['analysis_report'])
