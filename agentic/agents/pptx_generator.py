"""
A CrewAI agent to generate a PowerPoint (.pptx) file from structured Markdown content.

This module demonstrates how to:
1. Receive analyzed/structured Markdown data from the MarkdownAnalyzerAgent (JSON-like structure).
2. Generate slides in a PowerPoint file with text, bullet points, code listings, and basic images.
3. Compose a cohesive presentation while retaining the original Markdown structure.

Dependencies:
- python-pptx
- crewai

Example usage (see the test function at the end):
    from pptx_generator import PPTXGeneratorAgent
    generator = PPTXGeneratorAgent()
    pptx_file_path = generator.generate_pptx_from_analysis(conversion_plan_json, "output.pptx")
"""

import json
import os
import sys
import traceback
from pathlib import Path
from typing import Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

from crewai import Agent, Task, Crew
from crewai.tools import BaseTool


class PPTXBuilderTool(BaseTool):
    """
    Tool that takes a slide plan (JSON) and builds a PowerPoint file.
    """
    name: str = "BuildPPTX"
    description: str = "Constructs a PowerPoint file from a structured slide plan"

    def _run(self, slide_plan_json: str, output_pptx_path: str = "presentation.pptx") -> str:
        """
        Create a PowerPoint presentation (.pptx) from the provided slide plan JSON.
        
        Args:
            slide_plan_json: JSON string with the structured slides data.
            output_pptx_path: Where to save the resulting PPTX. Defaults to 'presentation.pptx'.
        
        Returns:
            File path of the created PPTX.
        """
        try:
            print(f"Building presentation at: {output_pptx_path}")
            slide_plan = json.loads(slide_plan_json)
            prs = Presentation()

            # Ensure the slides key exists
            if "slides" not in slide_plan and isinstance(slide_plan, dict):
                # Try to adapt the structure if it doesn't have a slides key but looks like a slide plan
                if any(key in slide_plan for key in ["total_slides", "has_title_slide"]):
                    print("Using slide plan as-is")
                else:
                    # Wrap the dictionary in a slides list if it seems to be a single slide
                    print("Wrapping input in slides key")
                    slide_plan = {"slides": [slide_plan]}

            for i, slide_info in enumerate(slide_plan.get("slides", [])):
                if i == 0 and slide_plan.get("has_title_slide", True):  # Default to True for has_title_slide
                    layout = prs.slide_layouts[0]  # Title slide layout
                else:
                    # Use a generic title & content layout
                    layout = prs.slide_layouts[1]

                slide = prs.slides.add_slide(layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]

                # Title
                title_text = slide_info.get("title", "Untitled Slide")
                title_shape.text = title_text

                # Convert bullet points, code blocks, and images in 'content'
                tf = body_shape.text_frame
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.clear()

                for content_item in slide_info.get("content", []):
                    ctype = content_item.get("type")
                    if ctype == "list":
                        for item in content_item.get("items", []):
                            p = tf.add_paragraph()
                            p.text = item
                            p.level = 0
                    elif ctype == "code":
                        code_lines = content_item.get("code", "").split("\n")
                        for line in code_lines:
                            p = tf.add_paragraph()
                            p.text = line
                            p.level = 1
                    elif ctype == "image":
                        # Insert a placeholder image (image path or basic shape).
                        image_path = content_item.get("src")
                        if image_path and os.path.exists(image_path):
                            left = top = Inches(2)
                            slide.shapes.add_picture(image_path, left, top, width=Inches(3))
                        else:
                            p = tf.add_paragraph()
                            p.text = "[Missing Image]"
                            p.level = 0
                    elif ctype == "subheading":
                        p = tf.add_paragraph()
                        p.text = content_item.get("text", "")
                        p.level = 0
                        p.font.bold = True
                    else:
                        # Fallback: treat as normal text
                        p = tf.add_paragraph()
                        p.text = content_item.get("text", "")
                        p.level = 0

            prs.save(output_pptx_path)
            return output_pptx_path
        except Exception as e:
            print(f"Error building PPTX: {str(e)}")
            import traceback
            traceback.print_exc()
            return f"Failed to build PPTX: {str(e)}"


class PPTXGeneratorAgent:
    """
    Agent that uses CrewAI to coordinate the creation of PowerPoint files
    from analyzed Markdown (JSON) structures.
    """

    def __init__(self):
        self.build_tool = PPTXBuilderTool()
        self.agent = self._setup_agent()

    def _setup_agent(self) -> Agent:
        """
        Creates a CrewAI Agent with the PPTXBuilderTool.
        """
        return Agent(
            role="PPTX Generation Agent",
            goal="Convert structured slide JSON into a PowerPoint file",
            backstory="Expert in building presentations from structured data.",
            tools=[self.build_tool],
            verbose=False,
            allow_delegation=False
        )

    def generate_pptx_from_analysis(self, slide_plan: dict, output_path: str) -> str:
        """
        High-level method to build the PPTX directly from a Python dict (slide_plan).
        
        Args:
            slide_plan: Dictionary containing the structured slides info (from the MarkdownAnalyzerAgent).
            output_path: Path to output the PPTX file.
        
        Returns:
            The filepath of the generated PPTX or an error message.
        """
        try:
            slide_plan_json = json.dumps(slide_plan)
            task_description = f"Generate PPTX from slide plan JSON. Output path: {output_path}"

            task = Task(
                description=task_description,
                agent=self.agent,
                expected_output="Path to created PPTX file"
            )

            # Create a single-step crew
            crew = Crew(agents=[self.agent], tasks=[task], verbose=False)

            # Kick off the build process
            result = self.build_tool._run(slide_plan_json, output_path)
            return result

        except Exception as e:
            return str(e)


def test_pptx_generator():
    """
    Demonstration of using the PPTXGeneratorAgent with sample data
    that might come from the MarkdownAnalyzerAgent.
    """
    # Add the project root to the Python path
    project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
    if project_root not in sys.path:
        sys.path.append(project_root)

    # Sample structured data simulating the output from the MarkdownAnalyzerAgent
    sample_conversion_plan = {
        "total_slides": 2,
        "has_title_slide": True,
        "slides": [
            {
                "title": "Title Slide",
                "content": [
                    {"type": "subheading", "text": "Presenter: Jane Doe"},
                    {"type": "subheading", "text": "Date: 2025-01-01"}
                ]
            },
            {
                "title": "First Content Slide",
                "content": [
                    {
                        "type": "list",
                        "list_type": "unordered",
                        "items": ["Introduction", "Main point", "Conclusion"]
                    },
                    {
                        "type": "code",
                        "language": "python",
                        "code": "def hello_world():\n    print('Hello, PPTX!')"
                    }
                ]
            }
        ]
    }

    generator = PPTXGeneratorAgent()
    output_file = generator.generate_pptx_from_analysis(sample_conversion_plan, "demo_output.pptx")
    print(f"Created PPTX: {output_file}")


if __name__ == "__main__":
    test_pptx_generator()
