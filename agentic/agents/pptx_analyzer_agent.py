from crewai import Agent
from crewai.tools import BaseTool
from agentic.agents.pptx_analyzer import PptxAnalyzerAgent
from pptx import Presentation

class AnalyzePresentation(BaseTool):
    """Tool for analyzing PowerPoint presentations."""
    
    name: str = "analyze_presentation"
    description: str = "Analyze a PowerPoint presentation and provide recommendations for improvement"
    
    def __init__(self, analyzer_func):
        super().__init__()
        self._analyzer_func = analyzer_func
    
    def _run(self, file_path: str, enhancement_options: dict = None) -> str:
        return self._analyzer_func(file_path, enhancement_options)

class PptxAnalyzerAgentWithAPIKey:
    """
    Extended version of PptxAnalyzerAgent that accepts api_key parameter
    for compatibility with other agents in the crew.
    """
    
    def __init__(self, api_key=None):
        """Initialize the analyzer agent with API key for compatibility."""
        self.api_key = api_key
        
        # Create the analyzer tool directly without relying on base_analyzer
        analyze_tool = AnalyzePresentation(analyzer_func=self.analyze_presentation)
        
        self.agent = Agent(
            role="PowerPoint Presentation Analyst",
            goal="Analyze PowerPoint presentations and provide detailed improvement recommendations",
            backstory=(
                "You are an expert presentation analyst with years of experience "
                "reviewing and critiquing PowerPoint presentations. You specialize "
                "in identifying design flaws, content issues, and structural problems "
                "in presentations."
            ),
            verbose=True,
            allow_delegation=False,
            tools=[analyze_tool]
        )
    
    def analyze_presentation(self, file_path, enhancement_options=None):
        """Analyze a PowerPoint presentation and provide recommendations."""
        try:
            # Perform our own analysis instead of relying on base_analyzer.analyze
            prs = Presentation(file_path)
            
            # Count slides and analyze content
            slide_count = len(prs.slides)
            has_title_slide = False
            text_heavy_slides = 0
            image_count = 0
            
            for slide in prs.slides:
                # Check if this is a title slide
                if slide.slide_layout.name.lower().startswith('title'):
                    has_title_slide = True
                
                # Count shapes with text
                text_shape_count = 0
                for shape in slide.shapes:
                    if hasattr(shape, "text") and len(shape.text) > 0:
                        text_shape_count += 1
                    if hasattr(shape, "image"):
                        image_count += 1
                
                # Consider slides with more than 3 text elements as text-heavy
                if text_shape_count > 3:
                    text_heavy_slides += 1
            
            # Format the recommendations based on analysis and enhancement options
            formatted_analysis = f"Analysis of presentation at {file_path}:\n\n"
            formatted_analysis += f"- Total slides: {slide_count}\n"
            formatted_analysis += f"- Has title slide: {'Yes' if has_title_slide else 'No'}\n"
            formatted_analysis += f"- Text-heavy slides: {text_heavy_slides} ({int(text_heavy_slides/slide_count*100)}% of total)\n"
            formatted_analysis += f"- Images: {image_count}\n\n"
            
            if enhancement_options and enhancement_options.get("improve_design", True):
                formatted_analysis += "Design recommendations:\n"
                formatted_analysis += "- Improve slide layout consistency\n"
                formatted_analysis += "- Enhance color scheme cohesion\n"
                formatted_analysis += "- Optimize typography for better readability\n\n"
                
            if enhancement_options and enhancement_options.get("improve_content", True):
                formatted_analysis += "Content recommendations:\n"
                formatted_analysis += "- Simplify complex text passages\n"
                formatted_analysis += "- Add more visual elements to support key points\n"
                formatted_analysis += "- Ensure consistent messaging throughout\n\n"
                
            if enhancement_options and enhancement_options.get("improve_structure", True):
                formatted_analysis += "Structure recommendations:\n"
                formatted_analysis += "- Improve slide transitions and flow\n"
                formatted_analysis += "- Reorganize content for better narrative\n"
                formatted_analysis += "- Add clear section markers\n\n"
                
            return formatted_analysis
            
        except Exception as e:
            return f"Error analyzing presentation: {str(e)}"
