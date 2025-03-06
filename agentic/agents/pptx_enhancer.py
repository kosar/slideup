import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR  # Ensure this import is present at the top

class PPTXEnhancer:
    def __init__(self, input_file, recommendations):
        # Load the existing PPTX and store recommendations.
        self.input_file = input_file
        self.recommendations = recommendations  # e.g. { "design": "Improve margins and alignment", "text": "Simplify text and improve clarity", "visual": "Add consistent imagery and icons" }
        self.prs = Presentation(input_file)

    def improve_slide_design(self, slide):
        # Skip layout changes for title slides (assumed to have a title placeholder)
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                try:
                    # Assuming type 1 is title placeholder
                    if shape.placeholder_format.type == 1:
                        return slide
                except Exception:
                    continue
        # Enhance slide design and layout based on recommendations.
        # For example, if design recommendation exists, shift shapes to a left margin of 0.5 inches.
        design_rec = self.recommendations.get("design", "")
        if "margins" in design_rec.lower():
            for shape in slide.shapes:
                if hasattr(shape, 'left'):
                    # Adjust left position to a fixed margin if not already close
                    if shape.left > Inches(0.5):
                        shape.left = Inches(0.5)
        # Adjust text shapes to use a uniform margin and width
        slide_width = self.prs.slide_width
        margin = Inches(0.5)
        # Reserve extra width if visual recommendation exists
        visual_rec = self.recommendations.get("visual", "")
        if "imagery" in visual_rec.lower():
            available_width = slide_width - Inches(3) - 2 * margin
        else:
            available_width = slide_width - 2 * margin
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    # Skip title placeholder
                    try:
                        if shape.placeholder_format.type == 1:
                            continue
                    except Exception:
                        pass
                shape.left = margin
                shape.width = available_width
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                # Example: adjust left margin if needed.
                pass  # Placeholder for improved layout logic.
        return slide

    def enhance_text(self, slide):
        # Improve text clarity and conciseness.
        text_rec = self.recommendations.get("text", "")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # If text is long, reduce font size for better clarity.
                    if len(paragraph.text) > 100 and "simplify" in text_rec.lower():
                        paragraph.font.size = Pt(16)
                    else:
                        paragraph.font.size = Pt(18)
        if hasattr(slide, "shapes"):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # Example: set a unified font size.
                        paragraph.font.size = Pt(18)
        return slide

    def add_visual_elements(self, slide, layout_toggle):
        # Skip image insertion on title slides to avoid overlap.
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                try:
                    if shape.placeholder_format.type == 1:
                        return slide
                except Exception:
                    continue
        visual_rec = self.recommendations.get("visual", "")
        pict_found = any(hasattr(shape, 'image') for shape in slide.shapes)
        if not pict_found and "imagery" in visual_rec.lower():
            placeholder_path = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'placeholder.png')
            try:
                margin = Inches(0.5)
                image_width = Inches(3)
                image_top = Inches(1.5)
                # Use layout_toggle to choose image side:
                if layout_toggle:
                    image_left = self.prs.slide_width - image_width - margin
                else:
                    image_left = margin
                slide.shapes.add_picture(placeholder_path, image_left, image_top, width=image_width)
            except Exception as e:
                print(f"Error adding visual element: {e}")
        return slide

    def ensure_consistency(self):
        # Enforce consistency across slides (e.g. uniform title styles).
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and 'title' in shape.name.lower():
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = "Calibri"
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(28)
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and 'title' in shape.name.lower():
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(28)
        return self.prs

    def enhance(self, output_file):
        layout_toggle = True
        for i, slide in enumerate(self.prs.slides):
            if i == 0:
                slide = self.improve_slide_design(slide)
                slide = self.enhance_text(slide)
                # Remove image shapes from the title slide, if any
                for shape in list(slide.shapes):
                    if hasattr(shape, 'image'):
                        try:
                            slide.shapes._spTree.remove(shape._element)
                        except Exception:
                            continue
                try:
                    # Explicitly adjust the title and subtitle placeholders.
                    title_shape = slide.shapes.title
                    subtitle_shape = slide.placeholders[1]
                    
                    # Set title placeholder: position near the top.
                    title_shape.top = Inches(1.25)
                    title_shape.left = Inches(1)
                    title_shape.width = self.prs.slide_width - Inches(2)
                    title_shape.height = Inches(1.2)
                    title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    for paragraph in title_shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(36)
                    
                    # Set subtitle placeholder (Subtitle, Presenter, Date) below the title.
                    subtitle_shape.top = Inches(3.25)
                    subtitle_shape.left = Inches(1)
                    subtitle_shape.width = self.prs.slide_width - Inches(2)
                    subtitle_shape.height = Inches(1.5)
                    subtitle_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    for paragraph in subtitle_shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(24)
                except Exception as e:
                    print("Error reformatting title slide:", e)
                    import sys
                    sys.exit(1)  # Exit instead of declaring success
                continue
            else:
                slide = self.improve_slide_design(slide)
                slide = self.enhance_text(slide)
                is_title = False
                for shape in slide.shapes:
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        try:
                            if shape.placeholder_format.type == 1:
                                is_title = True
                                break
                        except Exception:
                            continue
                if is_title:
                    for shape in list(slide.shapes):
                        if hasattr(shape, 'image'):
                            try:
                                if "placeholder.png" in shape.image.filename:
                                    slide.shapes._spTree.remove(shape._element)
                            except Exception:
                                continue
                else:
                    slide = self.add_visual_elements(slide, layout_toggle)
                    layout_toggle = not layout_toggle
        self.ensure_consistency()
        self.prs.save(output_file)
        print(f"Enhanced PPTX saved as {output_file}")

def main():
    import sys
    if len(sys.argv) < 3:
        print("Usage: python pptx_enhancer.py <input_pptx> <output_pptx>")
        sys.exit(1)
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    # Dummy recommendations; in practice these come from user input or a CrewAI workflow.
    recommendations = {
        "design": "Improve margins and alignment",
        "text": "Simplify text and improve clarity",
        "visual": "Add consistent imagery and icons and spacing on the page"
    }
    enhancer = PPTXEnhancer(input_file, recommendations)
    enhancer.enhance(output_file)

if __name__ == '__main__':
    main()
