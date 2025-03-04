import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import sys
import requests
import urllib.parse
import openai  # New import for OpenAI
from openai import OpenAI  # New import for OpenAI
import os  # New import for environment variables
import uuid  # New import for unique filenames
import time  # New import for delay
import base64
from stability_sdk import client
import stability_sdk.interfaces.gooseai.generation.generation_pb2 as generation
import io
import requests

# Define the directory to store all generated and temporary files
GENERATED_FILES_DIR = os.path.join(os.getcwd(), 'generated_files')

# Ensure the directory exists
if not os.path.exists(GENERATED_FILES_DIR):
    os.makedirs(GENERATED_FILES_DIR)

# Load and parse the markdown file
# Returns: List of slide data dictionaries
def parse_markdown(file_path):
    with open(file_path, 'r') as f:
        content = f.read()

    # Split slides with flexible header matching
    slides = re.split(r'^### Slide\s*\d+:\s*', content, flags=re.MULTILINE)[1:]
    slide_data = []
    
    for slide in slides:
        lines = [line.rstrip() for line in slide.split('\n') if line.strip()]
        if not lines:
            continue

        first_line = lines[0]
        
        # Title slide detection with case insensitivity
        if 'title slide' in first_line.lower():
            slide_info = {'type': 'title', 'content': {}}
            current_key = None
            
            for line in lines[1:]:
                # Handle different bullet styles and whitespace
                line = re.sub(r'^[*\-]\s*', '', line).strip()
                
                # Flexible key-value parsing
                match = re.match(r'\s*\**([^:]+?)\s*:\s*\**(.+)', line, re.I)
                if match:
                    current_key = match.group(1).strip().lower()
                    slide_info['content'][current_key] = match.group(2).strip()
                elif current_key:
                    # Append to existing key with line break
                    slide_info['content'][current_key] += '\n' + line
            
            # Ensure required fields exist with empty defaults
            for field in ['title', 'subtitle', 'presenter', 'date']:
                slide_info['content'].setdefault(field, '')
            
            slide_data.append(slide_info)
            continue

        # Process content slides
        current_slide = {
            'type': 'content',
            'title': first_line,
            'content': [],
            'references': []
        }

        collecting_references = False
        
        for line in lines[1:]:
            # Original reference detection that works
            if re.match(r'\* \*\*References:\*\*', line):
                collecting_references = True
                continue
                
            if collecting_references:
                # Original reference parsing that worked
                ref_match = re.match(r'\s*\* \[(.*)\]\((.*)\)', line)
                if ref_match:
                    current_slide['references'].append({
                        'text': ref_match.group(1).strip(),
                        'url': ref_match.group(2).strip()
                    })
                continue
            
            # Bullet parsing with bold preservation
            bullet_match = re.match(r'^(\s*)[*\-]\s+(.*)', line)
            if bullet_match:
                indent = len(bullet_match.group(1))
                text = bullet_match.group(2).strip()  # Preserves **bold** markers
                level = indent // 4  # 4 spaces per indent level
                
                current_slide['content'].append({
                    'text': text,
                    'level': min(level, 3)  # Limit to 3 levels
                })

        slide_data.append(current_slide)
    
    return slide_data

# Helper function to validate a URL and return a search URL if invalid
# Returns: Valid URL or search URL
# Used in the image generation process
def validate_url(url, title, authors):
    try:
        response = requests.head(url, allow_redirects=True, timeout=5)
        if response.status_code == 200:
            return url
        else:
            # Construct a search query URL
            query = f"{title} by {authors}"
            encoded_query = urllib.parse.quote_plus(query)
            search_url = f"https://www.google.com/search?q={encoded_query}"
            return search_url
    except requests.RequestException:
        # In case of request failure, return search query URL
        query = f"{title} by {authors}"
        encoded_query = urllib.parse.quote_plus(query)
        search_url = f"https://www.google.com/search?q={encoded_query}"
        return search_url

# Helper function to download and save the image
# Returns: Unique image filename or None if download fails
# Used in the image generation process
def download_image(image_url):
    try:
        image_response = requests.get(image_url)
        unique_image_name = f'temp_image_{uuid.uuid4()}.png'
        unique_image_path = os.path.join(GENERATED_FILES_DIR, unique_image_name)
        
        with open(unique_image_path, 'wb') as img_file:
            img_file.write(image_response.content)
            
        return unique_image_path  # Return the full path instead of just the name
        
    except Exception as e:
        print(f"Failed to download image: {e}")
        return None

# Helper function to clean the prompt text for DeepSeek image generation
# Returns: Cleaned text with technical terms removed
# Used in the image generation process
def clean_prompt(text):
    # Updated to use generate_clean_prompt with fallback
    return generate_clean_prompt(text)

# Helper function to extract theme keywords from a business concept
# This is useful because the DeepSeek API can provide theme keywords which can be used
# to generate images that are more relevant to the concept.
# Returns: List of theme keywords
# Used in the image generation process
def get_theme_keywords(concept):
    # Updated to use generate_theme_keywords with fallback
    return generate_theme_keywords(concept)

# Helper function to generate a prompt for DeepSeek image generation
# Returns: Prompt text for image generation or None if API key is missing
# Uses the DeepSeek API to generate a prompt for image generation
# This is useful because the DeepSeek API can provide prompts that are more
# tailored to the concept, leading to better image generation results
# The idea is to clean the prompt by removing technical terms and then
# generate a prompt that can be used to create an image
# Used in the image generation process
# TODO: This function can be further optimized by using a cache for repeated prompts
#       and by handling API errors more gracefully. 
#       In addition, the exact strategy and directive for what to clean may require generalization 
#       in case technical term filtering is not the preferred strategy. 
def generate_clean_prompt(text):
    api_key = os.getenv('DEEPSEEK_API_KEY')
    if not api_key:
        return _fallback_clean_prompt(text)

    try:
        response = generate_deepseek_output(text, prompt_type='clean_prompt')
        return response if response else _fallback_clean_prompt(text)
    except Exception as e:
        print(f"generate_clean_prompt failed: {e}")
        return _fallback_clean_prompt(text)

# Helper function to generate theme keywords for a business concept using Generative AI by default.
# Returns: Theme keywords or None if API key is missing
# This is useful because the DeepSeek API can provide theme keywords which can be used
# to generate images that are more relevant to the concept.
# Used in the image generation process
def generate_theme_keywords(concept):
    api_key = os.getenv('DEEPSEEK_API_KEY')
    if not api_key:
        return _fallback_get_theme_keywords(concept)

    try:
        response = generate_deepseek_output(concept, prompt_type='theme_keywords')
        return response if response else _fallback_get_theme_keywords(concept)
    except Exception as e:
        print(f"generate_theme_keywords failed: {e}")
        return _fallback_get_theme_keywords(concept)

# Helper function to create specific output from a specific prompt on content provided by the user. 
# This is used by other generator functions in this code to allow code reuse, with the prompt_type 
# governing what specifically the LLM is asked to output. 
# Returns: Generated output or None if API key is missing
# Used in the image generation process
# TODO: Opportunities to generalize this or combine it with other similar functions that call LLM APIs. 
#       Opportunities to leverage this function and its built-in system prompt flexibility to add robustness 
def generate_deepseek_output(input_text, prompt_type):
    try:
        client = OpenAI(
            api_key=os.getenv('DEEPSEEK_API_KEY'),
            base_url="https://api.deepseek.com/v1",
        )
        
        system_prompts = {
            'clean_prompt': "Clean the following text by removing technical terms and special characters:",
            'theme_keywords': "Extract up to three theme keywords related to the following business concept:",
        }
        
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": system_prompts[prompt_type]},
                {"role": "user", "content": input_text}
            ],
            temperature=0.7,
            max_tokens=100,
            top_p=0.9
        )
        
        return response.choices[0].message.content.strip()
        
    except Exception as e:
        print(f"DeepSeek API error: {e}")
        return None

# Internal function to provide fallback support if LLMs are offline
def _fallback_clean_prompt(text):
    # Original implementation of clean_prompt
    # Remove common technical terms that might trigger text generation
    technical_terms = [
        "API", "SDK", "UI", "database", "server", "code", "programming",
        "interface", "framework", "algorithm", "function", "variable"
    ]
    
    cleaned_text = text.lower()
    for term in technical_terms:
        cleaned_text = cleaned_text.replace(term.lower(), "")
    
    # Remove special characters and extra whitespace
    cleaned_text = re.sub(r'[^\w\s]', ' ', cleaned_text)
    cleaned_text = ' ' .join(cleaned_text.split())
    
    return cleaned_text

# Internal function to provide fallback support if LLMs are offline
def _fallback_get_theme_keywords(concept):
    # Original implementation of get_theme_keywords
    business_themes = {
        "growth": ["ascending curves", "upward movement", "organic growth", "expanding circles"],
        "innovation": ["light bulbs", "geometric patterns", "interconnected nodes", "abstract technology"],
        "strategy": ["chess pieces", "maze patterns", "compass", "pathways"],
        "success": ["mountain peaks", "ascending steps", "achievement symbols", "victory elements"],
        "collaboration": ["interlocking shapes", "connected circles", "bridge imagery", "unified elements"],
        "security": ["shield shapes", "protective layers", "lock patterns", "secure structures"],
        "efficiency": ["streamlined shapes", "flowing lines", "optimized patterns", "smooth transitions"],
        "development": ["building blocks", "foundation structures", "progressive elements", "staged growth"],
        "analysis": ["graph patterns", "data visualization", "analytical shapes", "measurement tools"],
        "leadership": ["guiding stars", "directional elements", "beacon imagery", "leading paths"]
    }
    
    concept_lower = concept.lower()
    selected_themes = []
    
    for theme, keywords in business_themes.items():
        if any(word in concept_lower for word in [theme] + keywords):
            selected_themes.extend(keywords)
    
    return selected_themes[:3]  # Return top 3 matching themes

# Helper function to extract theme keywords from a business concept
# Returns: List of theme keywords
# Used in the image generation process
def get_theme_keywords(concept):
    business_themes = {
        "growth": ["ascending curves", "upward movement", "organic growth", "expanding circles"],
        "innovation": ["light bulbs", "geometric patterns", "interconnected nodes", "abstract technology"],
        "strategy": ["chess pieces", "maze patterns", "compass", "pathways"],
        "success": ["mountain peaks", "ascending steps", "achievement symbols", "victory elements"],
        "collaboration": ["interlocking shapes", "connected circles", "bridge imagery", "unified elements"],
        "security": ["shield shapes", "protective layers", "lock patterns", "secure structures"],
        "efficiency": ["streamlined shapes", "flowing lines", "optimized patterns", "smooth transitions"],
        "development": ["building blocks", "foundation structures", "progressive elements", "staged growth"],
        "analysis": ["graph patterns", "data visualization", "analytical shapes", "measurement tools"],
        "leadership": ["guiding stars", "directional elements", "beacon imagery", "leading paths"]
    }
    
    concept_lower = concept.lower()
    selected_themes = []
    
    for theme, keywords in business_themes.items():
        if any(word in concept_lower for word in [theme] + keywords):
            selected_themes.extend(keywords)
    
    return selected_themes[:3]  # Return top 3 matching themes

# Helper function to generate a prompt for DeepSeek image generation
# Returns: Prompt text for image generation or None if API key is missing
# Used in the image generation process
def generate_deepseek_prompt(concept):
    api_key = os.getenv('DEEPSEEK_API_KEY')
    if not api_key:
        print("DeepSeek API key not found. Using fallback prompt generation.")
        return None

    try:
        client = OpenAI(
            api_key=api_key,
            base_url="https://api.deepseek.com/v1",  # DeepSeek's API endpoint
        )
        
        system_prompt = """You are a creative prompt engineer for AI image generation. Create VISUAL, ABSTRACT prompts for business concepts based on the content provided as context, with these rules:
                        1. Focus on symbolic representations, not literal
                        2. Use vivid color combinations and abstract shapes
                        3. Incorporate dynamic compositions
                        4. Reference artistic styles that resonate with a busines audience to complement words with visuals
                        5. Keep under 500 words
                        6. STYLE: Abstract, symbolic imagery. Use dynamic shapes, flowing lines, and vibrant colors to convey the essence of the content.
                        7. NO WORDS should be in the image unless they are well understood and spelled correctly. Emphasizes abstract visual elements. BLUR TEXT in your images if generated. When conveying scientific concepts, use metaphorical representations.
                        8. COMPOSITION: Create a balanced, harmonious composition. Use symmetry, contrast, and focal points to guide the viewer's eye."""

        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Here are the concepts for which we want to generate an image: {concept}\nGenerate a prompt that can be provided to an image creation model."}
            ],
            temperature=0.9,
            max_tokens=550,
            top_p=0.95
        )

        return response.choices[0].message.content.strip()
        
    except openai.error.InvalidRequestError as e:
        if 'Insufficient Balance' in str(e):
            print("DeepSeek API Error: Insufficient Balance. Using fallback prompt generation.")
            return None
        else:
            print(f"DeepSeek prompt generation failed: {e}")
            return None
    except Exception as e:
        print(f"DeepSeek prompt generation failed: {e}")
        return None

# Helper function to generate an image using the Stability API
# Returns: Unique image filename or None if generation fails
#  Note: the file  name is returned instead of the image data
# as the image data is not used in the final presentation
# and the file name is used to add the image to the presentation
# and then delete the temporary file
# The image data can be saved to a file and returned if needed
# in the future
def generate_image_stability(prompt, stability_prompt_override=None):
    api_key = os.getenv('STABILITY_API_KEY')
    if not api_key:
        print("Stability API key not found. Skipping image generation.")
        return None

    try:
        # Initialize Stability API client
        stability_api = client.StabilityInference(
            key=api_key,
            verbose=True,
        )
        
        # Get enhanced prompt from DeepSeek
        # deepseek_prompt = generate_deepseek_prompt(prompt)
        deepseek_prompt=None # temporarily stop generating prompts and use the static one

        if not deepseek_prompt:
            # Construct fallback prompt with a warning about DeepSeek failure
            print("Using fallback prompt due to DeepSeek prompt generation failure.")
            cleaned_concept = clean_prompt(prompt)
            theme_keywords = get_theme_keywords(cleaned_concept)   
            stability_prompt = stability_prompt_override or (
                f"Create a conceptual visualization of an abstract concept representing: {cleaned_concept} " 
                f"STYLE: Abstract, symbolic imagery. Use dynamic shapes, flowing lines, and vibrant colors to convey the essence of the content. " 
                f"Explore color symbolism and metaphorical representations. "
                f"Avoid realistic objects or scenes. Focus on conveying the underlying ideas. "
                f"DO NOT include text, numbers, or literal interpretations. Only abstract visual elements. BLUR TEXT in your images if generated. COMPOSITION: Create a balanced, harmonious composition. Use symmetry, contrast, and focal points to guide the viewer's eye. No sexuality or nudity or suggestive content. No text in the image. No violence or gore."
                f"Visual themes if any: {', '.join(theme_keywords)}"
            )
        else:
            # Use the DeepSeek-generated prompt
            print (f"Using DeepSeek prompt (len): {len(deepseek_prompt)}")
            stability_prompt = (
                f"Image generation instructions: {deepseek_prompt}"
                f"Emphasize abstract visual elements that represent the concept provided: {prompt}"
            )

        # pretty print the stability_prompt
        # print(f"stability_prompt:\n{stability_prompt}")

        # Enhanced parameters
        answers = stability_api.generate(
            prompt=stability_prompt,
            steps=60,
            cfg_scale=8.0,
            width=1024,
            height=1024,
            samples=1,
            sampler=generation.SAMPLER_K_DPMPP_2M,
            style_preset='conceptual',
        )

        # Process the generated image
        for resp in answers:
            for artifact in resp.artifacts:
                if artifact.finish_reason == generation.FILTER:
                    print("Content was filtered. Trying again with a modified prompt...")
                    continue
                    
                if artifact.type == generation.ARTIFACT_IMAGE:
                    img_data = io.BytesIO(artifact.binary)
                    unique_image_name = f'temp_image_{uuid.uuid4()}.png'
                    unique_image_path = os.path.join(GENERATED_FILES_DIR, unique_image_name)
                    
                    with open(unique_image_path, 'wb') as img_file:
                        img_file.write(img_data.getvalue())
                    
                    return unique_image_path

        return None
        
    except Exception as e:
        print(f"generate_image_stability Exception: {e}")
        return None

# Helper function to generate speaker notes for a slide content
# Returns: Speaker notes as plain text or an empty string if generation fails
# Used in the presentation creation process
def generate_speaker_notes(content, speaker_notes_prompt_override=None):
    api_key = os.getenv('OPENAI_API_KEY')
    if not api_key:
        print("OpenAI API key not found. Skipping speaker notes generation.")
        return ""

    try:
        client = OpenAI(api_key=api_key)
        
        # Structure the content for better context
        if isinstance(content, dict):
            formatted_content = (
                f"Title: {content.get('title', '')}\n"
                "Content:\n" + "\n".join(f"- {point['text']}" 
                for point in content.get('content', []) if point.get('text'))
            )
        else:
            formatted_content = str(content)

        # Create a detailed prompt for better speaker notes
        system_prompt = speaker_notes_prompt_override or """You are an expert presentation coach creating speaker notes for a slide. Provide clear, concise, and comprehensive talking points that expand on each bullet point without using markdown or any special formatting. Return the notes as plain text."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Create speaker notes for this slide content:\n\n{formatted_content}"}
        ]

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0.7,
            max_tokens=500,
            top_p=1.0,
            frequency_penalty=0.0,
            presence_penalty=0.0
        )

        notes = response.choices[0].message.content.strip()
        return notes

    except Exception as e:
        print(f"Exception caught: Failed to generate speaker notes: {e}")
        return ""

# Helper function to create a PowerPoint presentation from slide data
# Returns: None
# Saves the presentation to the specified output path
# Used in the main program execution
# Updated to handle custom slide layouts and image generation
# This is the heart of the presentation generation process
# It creates a PowerPoint presentation with custom slide layouts
# and dynamically adjusts font sizes based on content density
# It also generates speaker notes for each slide if requested
# and adds images to slides based on the content
# The image generation process uses the Stability API for image generation
# and the DeepSeek API for generating prompts for the Stability API
# The speaker notes generation uses the OpenAI API for text generation
# The process is optimized for speed and efficiency
# with dynamic font sizing and layout adjustments
# based on the content of each slide
# TODO: this may need to be refactored into smaller functions as it is getting complex
def create_presentation(slides, output_path, add_notes=False, add_images_stability=False, stability_prompt=None, speaker_notes_prompt=None):
    prs = Presentation()
    
    # Set slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    layout_toggle = True  # Initialize layout toggle

    total_slides = len(slides)
    index = 1
    for slide in slides:
        print(f"Processing slide {index} of {total_slides}")
        start_time = time.time()
        
        if slide['type'] == 'title':
            # Title slide layout remains unchanged
            slide_layout = prs.slide_layouts[0]
            sld = prs.slides.add_slide(slide_layout)
            title = sld.shapes.title
            subtitle = sld.placeholders[1]
            
            title.text = slide['content'].get('title', 'Title')
            subtitle.text = (f"{slide['content'].get('subtitle', '')}\n"
                           f"Presenter: {slide['content'].get('presenter', '')}\n"
                           f"Date: {slide['content'].get('date', '')}")
            
            title.text_frame.paragraphs[0].font.size = Pt(32)
            subtitle.text_frame.paragraphs[0].font.size = Pt(16)
        
        else:
            # Content slide with improved layout
            slide_layout = prs.slide_layouts[6]  # Blank slide for custom layout
            sld = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_box = sld.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = slide['title']
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            
            # Define fixed widths
            fixed_content_width = Inches(6)
            fixed_image_width = Inches(6)
            margin = Inches(0.5)

            # Determine positions based on layout toggle
            if add_images_stability:
                if layout_toggle:
                    # Image on the right, text on the left
                    content_left = margin
                    image_left = content_left + fixed_content_width + margin
                else:
                    # Image on the left, text on the right
                    image_left = margin
                    content_left = image_left + fixed_image_width + margin
                content_width = fixed_content_width
                image_width = fixed_image_width
            else:
                # Full width for content if no images
                content_left = margin
                content_width = Inches(12) - 2 * margin
                image_left = None  # No image

            # Add content textbox
            content_box = sld.shapes.add_textbox(
                content_left,           # Adjusted left position
                Inches(1.5),            # Top margin (unchanged)
                content_width,          # Adjusted width
                Inches(5)               # Height
            )
            
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            
            # Add content with proper formatting
            for point in slide['content']:
                if point['text']:
                    p = tf.add_paragraph()
                    p.level = point['level']
                    p.alignment = PP_ALIGN.LEFT
                    p.font.size = Pt(18)
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)
                    
                    # Handle bold formatting
                    text = point['text']
                    bold_split = re.split(r'(\*\*[^*]+\*\*)', text)
                    for segment in bold_split:
                        if segment.startswith('**') and segment.endswith('**'):
                            run = p.add_run()
                            run.text = segment.strip('**')
                            run.font.bold = True
                        else:
                            run = p.add_run()
                            run.text = segment

            # Calculate content metrics
            num_bullets = len(tf.paragraphs)
            total_chars = sum(len(para.text) for para in tf.paragraphs)
            
            # Dynamic font sizing based on content density and image presence
            if add_images_stability:
                # Tighter thresholds when images are present
                if num_bullets > 6 or total_chars > 600:  # 6 bullets or 600 chars
                    new_font_size = 16
                elif num_bullets > 4 or total_chars > 400:  # 4-6 bullets or 400-600 chars
                    new_font_size = 17
                else:
                    new_font_size = 18
            else:
                # Normal thresholds for full-width slides
                if num_bullets > 8 or total_chars > 800:  # 8 bullets or 800 chars
                    new_font_size = 16
                elif num_bullets > 5 or total_chars > 500:  # 5-8 bullets or 500-800 chars
                    new_font_size = 17
                else:
                    new_font_size = 18

            # Apply dynamic font sizing to all paragraphs and runs
            for para in tf.paragraphs:
                para.font.size = Pt(new_font_size)
                for run in para.runs:
                    run.font.size = Pt(new_font_size)
            
            # Add image if requested
            if add_images_stability:
                # Extract key points for image generation             
                key_points = ' '.join([point['text'] for point in slide['content']])
                if len(key_points.split()) < 3:  # Minimum content check
                    print(f"Skipping image for sparse content: {slide['title']}")
                    continue

                try:
                    # Generate image - now returns a local file path within GENERATED_FILES_DIR
                    image_path = generate_image_stability(key_points, stability_prompt)
                    
                    if image_path:
                        try:
                            # Add image to determined side
                            left = image_left
                            top = Inches(1.5)  # Align with content top
                            sld.shapes.add_picture(
                                image_path,
                                left,
                                top,
                                width=image_width,
                                height=image_width  # Keep aspect ratio 1:1
                            )
                            
                            # Clean up the temporary image file
                            os.remove(image_path)
                            
                        except Exception as e:
                            print(f"Failed to add image to slide: {e}")
                            # Clean up the temporary file in case of error
                            if os.path.exists(image_path):
                                os.remove(image_path)
                                
                except Exception as e:
                    print(f"Failed to generate image: {e}")

            # Add references if present
            if slide.get('references'):
                ref_top = Inches(6.5)  # Position at bottom of slide
                txBox = sld.shapes.add_textbox(Inches(0.5), ref_top, Inches(12), Inches(0.75))
                tf = txBox.text_frame
                
                p = tf.add_paragraph()
                p.text = 'References:'
                p.font.size = Pt(12)
                p.font.bold = True
                
                for ref in slide['references']:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = ref['text']
                    run.hyperlink.address = ref['url']
                    p.font.size = Pt(10)

            # Add slide number (new) to the bottom right corner
            slide_number = str(index)
            text_width = Inches(1)
            text_height = Inches(0.3)
            left = prs.slide_width - Inches(0.5) - text_width  # 0.5" from right edge
            top = prs.slide_height - Inches(0.3) - text_height  # 0.3" from bottom
            text_box = sld.shapes.add_textbox(left, top, text_width, text_height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = False
            p = text_frame.add_paragraph()
            p.text = slide_number
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(12)

            
            # Add speaker notes
            if add_notes:
                print(f"Generating speaker notes for slide titled: '{slide['title']}'")
                slide_content_text = "\n".join([point['text'] for point in slide['content']])
                speaker_notes = generate_speaker_notes(slide_content_text, speaker_notes_prompt)
                if speaker_notes:
                    notes_slide = sld.notes_slide
                    if notes_slide and notes_slide.notes_text_frame:
                        notes_slide.notes_text_frame.add_paragraph().text = speaker_notes
                    else:
                        print(f"Failed to add speaker notes for slide titled: '{slide['title']}'")
                else:
                    print(f"No speaker notes generated for slide titled: '{slide['title']}'")
    
            # Toggle layout for next slide
            layout_toggle = not layout_toggle

        elapsed_time = time.time() - start_time
        print(f"Finished slide {index} in {elapsed_time:.2f} seconds")
        index += 1

    prs.save(output_path)  # Save the presentation to the specified output path

# Main program execution
# Parses the input markdown file and generates a PowerPoint presentation
# Supports optional command-line arguments for speaker notes and image generation
# This code sets flags such as add_notes and add_images_stability which control
# the behavior of the presentation generation process
# The main function orchestrates the entire process
# by calling the necessary functions in sequence
# It also handles error messages and exceptions
# to provide a smooth user experience
# The program can be run from the command line
# with the markdown file as the first argument
# and optional flags for speaker notes and image generation
# The resulting PowerPoint presentation is saved to a file
# with the same name as the markdown file but with a .pptx extension
# The program execution time is displayed at the end
# Be sure to set the required environment variables for the APIs
# before running the program
# In particular, you must have valid API keys for OpenAI, Stability, and DeepSeek
# to use the full functionality of the program
# Typical Powerpoint slide generation is relatively inexpensive using this trio of APIs
# A typical presentation with 10-20 slides can be generated in a few minutes, and cost 
# about 10-20 cents for the image generation and text generation services
# The program can be run on any machine with Python installed
# and the required libraries installed via pip
# The resulting PowerPoint presentation can be viewed and edited
# in any PowerPoint software that supports the .pptx format
# One could use this program to generate presentations from markdown files and then upload / import
# the file directly into Google Slides. 
def main(cli_args=None, alt_output=None):
    if cli_args is None:
        cli_args = sys.argv

    # Remove empty strings from arguments (from optional flags not selected)
    cli_args = [arg for arg in cli_args if arg]

    if len(cli_args) < 2:
        print("Usage: python md2pptx.py <markdown_file> [--add-notes] [--add-images-stability] [--stability-prompt <prompt>] [--speaker-notes-prompt <prompt>]")
        sys.exit(1)

    markdown_file = cli_args[1]
    add_notes = '--add-notes' in cli_args
    add_images_stability = '--add-images-stability' in cli_args

    if alt_output:
        output_file = alt_output
    else:
        output_file = os.path.join(GENERATED_FILES_DIR, markdown_file.rsplit('.', 1)[0] + '.pptx')

    stability_prompt = None
    speaker_notes_prompt = None

    try:
        stability_prompt_index = cli_args.index("--stability-prompt") + 1
        stability_prompt = cli_args[stability_prompt_index]
    except (ValueError, IndexError):
        pass

    try:
        speaker_notes_prompt_index = cli_args.index("--speaker-notes-prompt") + 1
        speaker_notes_prompt = cli_args[speaker_notes_prompt_index]
    except (ValueError, IndexError):
        pass

    print(f"[DEBUG] Generating presentation from {markdown_file} with notes={add_notes}, images_stability={add_images_stability}, stability_prompt={stability_prompt}, speaker_notes_prompt={speaker_notes_prompt}")
    try:
        start_time = time.time()
        slides = parse_markdown(markdown_file)
        create_presentation(slides, output_file, add_notes, add_images_stability, stability_prompt, speaker_notes_prompt)
        end_time = time.time()
        elapsed_time = end_time - start_time
        print(f"Presentation created successfully from {markdown_file}. Saved as {output_file}")
        print(f"Total time taken: {elapsed_time:.2f} seconds")
    except FileNotFoundError:
        print(f"Error: File '{markdown_file}' not found.")
    except Exception as e:
        print(f"An error occurred: {e}")

if __name__ == '__main__':
    main()