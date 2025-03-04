import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from .md2pptx import generate_image_stability  # Import existing image generation function
except ImportError:
    from md2pptx import generate_image_stability

def extract_slide_content(prs):
    """
    Extract text content from each slide for image generation.
    Returns a list of dictionaries: [{'index': int, 'text': str}, ...]
    """
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text.strip())
        slides_text.append({
            'index': i,
            'text': " ".join(slide_text).strip()
        })
    return slides_text

def add_images_to_slides(prs, slides_text, adaptive_body_font_size):
    """
    For each slide, generate an image based on its text,
    then add the image in a layout similar to md2pptx.
    """
    layout_toggle = True
    placeholder_path = os.path.join(os.path.dirname(__file__), 'placeholder.png')  # Local fallback

    try:
        from openai.error import OpenAIError
    except ImportError:
        OpenAIError = Exception  # Fallback if openai.error is not available

    for slide_info in slides_text:
        slide_index = slide_info['index']
        # Skip the first slide (assumed title slide)
        if slide_index == 0:
            continue
        text_content = slide_info['text']
        if not text_content:
            continue  # Skip if no text

        # Generate image using the existing pipeline
        try:
            image_path = generate_image_stability(text_content)
            if not image_path:
                image_path = placeholder_path
        except OpenAIError as e:
            print(f"generate_image_stability Exception: {e}")
            image_path = placeholder_path

        sld = prs.slides[slide_index]
        margin = Inches(0.5)
        image_box_width = Inches(4)  # Image size remains the same
        slide_width = prs.slide_width
        # Calculate available width for body text
        body_text_width = slide_width - image_box_width - 2 * margin
        # Restrict the text block to a maximum width (e.g., 6 inches)
        max_text_width = Inches(6)
        text_width = min(body_text_width, max_text_width)
        # Move image upward (previous adjustment)
        image_top = Inches(2)
        
        # Separate title and body text shapes
        title_shape = None
        body_shapes = []
        for shape in sld.shapes:
            if shape.has_text_frame:
                if (shape.is_placeholder and shape.placeholder_format.type == 1) or ("title" in shape.name.lower()):
                    if title_shape is None:
                        title_shape = shape
                        continue
                body_shapes.append(shape)
        
        # Process title if found: position just above body text
        if title_shape is not None:
            title_shape.top = Inches(1.5)
            title_shape.left = margin
            title_shape.width = Inches(9)
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
        
        # Process body text shapes and set image_left based on layout
        if layout_toggle:
            # Image on right: position text on left side
            for shape in body_shapes:
                shape.top = Inches(2.4)
                shape.left = margin
                shape.width = text_width  # Use restricted text width
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                shape.text_frame.word_wrap = True
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(adaptive_body_font_size)
            image_left = slide_width - image_box_width - margin
        else:
            # Image on left: position text on right side
            for shape in body_shapes:
                shape.top = Inches(2.4)
                shape.left = margin + image_box_width + margin
                shape.width = text_width  # Use restricted text width
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                shape.text_frame.word_wrap = True
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(adaptive_body_font_size)
            image_left = margin

        # Insert picture on the chosen side with updated horizontal position
        try:
            sld.shapes.add_picture(
                image_path, image_left, image_top, width=image_box_width, height=image_box_width
            )
        except Exception as e:
            print(f"Error adding image on slide {slide_index}: {e}")
        finally:
            if image_path and image_path != placeholder_path and os.path.exists(image_path):
                os.remove(image_path)
        
        layout_toggle = not layout_toggle

def enhance_pptx(input_file, output_file=None):
    """
    Main function: extracts text from existing PPTX, generates images,
    places them on each slide using a layout toggle approach,
    then saves the enhanced PPTX.
    """
    if not output_file:
        base_name = os.path.splitext(input_file)[0]
        output_file = f"{base_name}_enhanced.pptx"
    
    prs = Presentation(input_file)
    slides_text = extract_slide_content(prs)
    # Compute adaptive font size based on longest text across slides
    max_text_len = max(len(slide['text']) for slide in slides_text) if slides_text else 0
    if max_text_len <= 200:
        adaptive_font_size = 24
    elif max_text_len <= 400:
        # Scale linearly from 24 (for 200 chars) down to 18 (for 400 chars)
        adaptive_font_size = round(24 - (max_text_len - 200) * (6 / 200))
    else:
        # For very long text, reduce further but not below 14pt
        adaptive_font_size = max(14, round(18 - (max_text_len - 400) * (4 / 200)))
    add_images_to_slides(prs, slides_text, adaptive_font_size)
    prs.save(output_file)
    print(f"Enhanced PPTX saved as {output_file}")

def main(cli_args=None):
    if cli_args is None:
        cli_args = sys.argv

    if len(cli_args) < 2:
        print("Usage: python pptx2enhancedpptx.py <existing_pptx_file> [output_file]")
        sys.exit(1)

    input_file = cli_args[1]
    output_file = cli_args[2] if len(cli_args) > 2 else None
    enhance_pptx(input_file, output_file)

if __name__ == '__main__':
    main()